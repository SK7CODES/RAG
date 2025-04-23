import os
import mimetypes
from typing import List, Dict, Any, Optional, Union
import tempfile
import json

# For PDF processing
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    pass

# For image processing
try:
    from PIL import Image
except ImportError:
    pass

# For audio processing
try:
    from pydub import AudioSegment
except ImportError:
    pass

# For video processing
try:
    import cv2
    import moviepy.editor as mp
except ImportError:
    pass

# For document processing
try:
    import docx2txt
    from pptx import Presentation
except ImportError:
    pass

class FileProcessor:
    """
    Processes different types of files for the multimodal RAG system.
    """
    
    def __init__(self):
        """Initialize the file processor."""
        self.supported_extensions = {
            'pdf': self._process_pdf,
            'txt': self._process_text,
            'png': self._process_image,
            'jpg': self._process_image,
            'jpeg': self._process_image,
            'mp3': self._process_audio,
            'wav': self._process_audio,
            'mp4': self._process_video,
            'docx': self._process_docx,
            'pptx': self._process_pptx
        }
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        Process a file based on its extension.
        
        Args:
            file_path (str): Path to the file to process
            
        Returns:
            Dict[str, Any]: Dict containing processed content and metadata
        """
        try:
            # Get file extension
            _, ext = os.path.splitext(file_path)
            ext = ext.lstrip('.').lower()
            
            # Check if we support this file type
            if ext not in self.supported_extensions:
                return {
                    'success': False,
                    'error': f"Unsupported file type: {ext}",
                    'file_path': file_path
                }
            
            # Process the file
            processor_func = self.supported_extensions[ext]
            result = processor_func(file_path)
            
            # Add file path to result
            result['file_path'] = file_path
            result['file_type'] = ext
            
            return result
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'file_path': file_path
            }
    
    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process PDF files and extract text content."""
        try:
            # Use PDFPlumber for more accurate text extraction
            text_content = []
            
            with pdfplumber.open(file_path) as pdf:
                metadata = {
                    'pages': len(pdf.pages),
                    'filename': os.path.basename(file_path)
                }
                
                # Extract text from each page
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"--- Page {i+1} ---\n{page_text}")
            
            # Alternative extraction with PyPDF2 if needed
            if not text_content:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    metadata = {
                        'pages': len(reader.pages),
                        'filename': os.path.basename(file_path)
                    }
                    
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text_content.append(f"--- Page {i+1} ---\n{page_text}")
            
            return {
                'success': True,
                'content_type': 'text',
                'content': '\n\n'.join(text_content),
                'metadata': metadata
            }
        except Exception as e:
            return {
                'success': False,
                'error': f"Error processing PDF: {str(e)}"
            }
    
    def _process_text(self, file_path: str) -> Dict[str, Any]:
        """Process text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            metadata = {
                'filename': os.path.basename(file_path),
                'size': os.path.getsize(file_path)
            }
            
            return {
                'success': True,
                'content_type': 'text',
                'content': content,
                'metadata': metadata
            }
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    content = file.read()
                
                metadata = {
                    'filename': os.path.basename(file_path),
                    'size': os.path.getsize(file_path)
                }
                
                return {
                    'success': True,
                    'content_type': 'text',
                    'content': content,
                    'metadata': metadata
                }
            except Exception as e:
                return {
                    'success': False,
                    'error': f"Error processing text file: {str(e)}"
                }
        except Exception as e:
            return {
                'success': False,
                'error': f"Error processing text file: {str(e)}"
            }
    
    def _process_image(self, file_path: str) -> Dict[str, Any]:
        """Process image files and extract metadata."""
        try:
            image = Image.open(file_path)
            
            metadata = {
                'filename': os.path.basename(file_path),
                'format': image.format,
                'mode': image.mode,
                'width': image.width,
                'height': image.height,
                'size': os.path.getsize(file_path)
            }
            
            # For RAG, we don't store the actual image content
            # but keep the path to the image file
            
            return {
                'success': True,
                'content_type': 'image',
                'content': f"Image file: {os.path.basename(file_path)}",
                'binary_path': file_path,
                'metadata': metadata
            }
        except Exception as e:
            return {
                'success': False,
                'error': f"Error processing image: {str(e)}"
            }
    
    def _process_audio(self, file_path: str) -> Dict[str, Any]:
        """Process audio files and extract metadata."""
        try:
            audio = AudioSegment.from_file(file_path)
            
            metadata = {
                'filename': os.path.basename(file_path),
                'duration_seconds': len(audio) / 1000,
                'channels': audio.channels,
                'sample_width': audio.sample_width,
                'frame_rate': audio.frame_rate,
                'size': os.path.getsize(file_path)
            }
            
            # For RAG, we don't store the actual audio content
            # but keep the path to the audio file
            
            return {
                'success': True,
                'content_type': 'audio',
                'content': f"Audio file: {os.path.basename(file_path)}",
                'binary_path': file_path,
                'metadata': metadata
            }
        except Exception as e:
            return {
                'success': False,
                'error': f"Error processing audio: {str(e)}"
            }
    
    def _process_video(self, file_path: str) -> Dict[str, Any]:
        """Process video files and extract metadata."""
        try:
            # Open the video file
            video = cv2.VideoCapture(file_path)
            
            # Extract metadata
            frame_count = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = video.get(cv2.CAP_PROP_FPS)
            width = int(video.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(video.get(cv2.CAP_PROP_FRAME_HEIGHT))
            
            # Release the video resource
            video.release()
            
            # Calculate duration
            duration_seconds = frame_count / fps if fps > 0 else 0
            
            metadata = {
                'filename': os.path.basename(file_path),
                'duration_seconds': duration_seconds,
                'frame_count': frame_count,
                'fps': fps,
                'width': width,
                'height': height,
                'size': os.path.getsize(file_path)
            }
            
            # For RAG, we don't store the actual video content
            # but keep the path to the video file
            
            return {
                'success': True,
                'content_type': 'video',
                'content': f"Video file: {os.path.basename(file_path)}",
                'binary_path': file_path,
                'metadata': metadata
            }
        except Exception as e:
            return {
                'success': False,
                'error': f"Error processing video: {str(e)}"
            }
    
    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Process DOCX files and extract text content."""
        try:
            # Extract text from DOCX file
            text = docx2txt.process(file_path)
            
            metadata = {
                'filename': os.path.basename(file_path),
                'size': os.path.getsize(file_path)
            }
            
            return {
                'success': True,
                'content_type': 'text',
                'content': text,
                'metadata': metadata
            }
        except Exception as e:
            return {
                'success': False,
                'error': f"Error processing DOCX: {str(e)}"
            }
    
    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process PPTX files and extract text content."""
        try:
            prs = Presentation(file_path)
            
            # Extract text from slides
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
            
            metadata = {
                'filename': os.path.basename(file_path),
                'slides': len(prs.slides),
                'size': os.path.getsize(file_path)
            }
            
            return {
                'success': True,
                'content_type': 'text',
                'content': '\n\n'.join(text_content),
                'metadata': metadata
            }
        except Exception as e:
            return {
                'success': False,
                'error': f"Error processing PPTX: {str(e)}"
            }
    
    def get_mime_type(self, file_path: str) -> str:
        """Get the MIME type of a file."""
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type:
            return mime_type
        
        # Fallback based on extension
        ext = os.path.splitext(file_path)[1].lower()
        if ext in ['.jpg', '.jpeg']:
            return 'image/jpeg'
        elif ext == '.png':
            return 'image/png'
        elif ext == '.pdf':
            return 'application/pdf'
        elif ext == '.txt':
            return 'text/plain'
        elif ext == '.mp3':
            return 'audio/mpeg'
        elif ext == '.wav':
            return 'audio/wav'
        elif ext == '.mp4':
            return 'video/mp4'
        elif ext == '.docx':
            return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif ext == '.pptx':
            return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        else:
            return 'application/octet-stream' 